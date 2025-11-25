from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_slide(prs, title, content_points):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Content
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for i, point in enumerate(content_points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(18)
        p.space_after = Pt(10)

def generate_presentation(output_file="presentation.pptx"):
    prs = Presentation()
    
    # --- Slide 1: Capa ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    title.text = "Detecção de Memes Enganosos"
    subtitle.text = "PBL 4 - Engenharia de Computação\nArquitetura Multimodal e Pipeline de IA"

    # --- Slide 2: Problema ---
    create_slide(prs, "O Problema: Desinformação Visual", [
        "Memes são vetores rápidos de Fake News.",
        "Dificuldade de moderação automática devido à multimodalidade (texto + imagem).",
        "Manipulações sutis (troca de contexto, edição de números).",
        "Objetivo: Criar um classificador automático robusto e explicável."
    ])

    # --- Slide 3: Metodologia Geral ---
    create_slide(prs, "Metodologia: Pipeline Multimodal", [
        "Abordagem de Fusão Tardia (Late Fusion).",
        "Processamento paralelo de Visão e Linguagem.",
        "Validação Semântica Externa (LLM).",
        "Foco em explicabilidade (XAI) e integração via API."
    ])

    # --- Slide 4: Dataset Sintético ---
    create_slide(prs, "Dataset e Geração de Dados", [
        "Escassez de dados rotulados em PT-BR.",
        "Solução: Gerador Sintético (Python/Pillow).",
        "Classes: Autêntico vs. Manipulado.",
        "Manipulações simuladas: Splicing de texto, métricas absurdas, degradação JPEG, inserção de ícones."
    ])

    # --- Slide 5: Pré-processamento & OCR ---
    create_slide(prs, "Pré-processamento e OCR", [
        "Tratamento de Imagem: Denoise, Equalização de Histograma.",
        "OCR: Tesseract 5.0 com configuração customizada (--psm 6).",
        "Extração de Metadados: Confiança média das palavras e bounding boxes.",
        "Detecção de anomalias no texto baseada em baixa confiança do OCR."
    ])

    # --- Slide 6: Modelo Visual (CNN) ---
    create_slide(prs, "Componente Visual", [
        "Backbone: MobileNetV2 (pré-treinado no ImageNet).",
        "Feature Extraction: Camada 'features' + Global Avg Pooling.",
        "Vetor de saída: 1280 dimensões.",
        "Vantagem: Leve e eficiente para inferência em CPU/Docker."
    ])

    # --- Slide 7: Modelo Textual (NLP) ---
    create_slide(prs, "Componente Textual", [
        "Tokenização customizada.",
        "Arquitetura: Embedding + Bidirectional LSTM.",
        "Captura de contexto sequencial (ex: 'não' antes de um verbo).",
        "Vetor de saída: 256 dimensões."
    ])

    # --- Slide 8: Fusão e Classificação ---
    create_slide(prs, "Modelo de Fusão", [
        "Concatenação: [Visual (1280) | Textual (256) | OCR Stats (3)].",
        "MLP Classifier: Dense -> Batch Norm -> ReLU -> Dropout.",
        "Saída: Probabilidade de ser 'Enganoso'.",
        "Integração LLM: Análise semântica pós-processamento (Gemini) para validação lógica."
    ])

    # --- Slide 9: Resultados Preliminares ---
    create_slide(prs, "Resultados e Métricas", [
        "Avaliação em conjunto de validação sintético.",
        "Métricas Chave: Acurácia, F1-Score e Matriz de Confusão.",
        "Explicabilidade Visual: Heatmaps via Grad-CAM.",
        "Transparência: API retorna evidências textuais e visuais."
    ])

    # --- Slide 10: Conclusão ---
    create_slide(prs, "Conclusões e Próximos Passos", [
        "Pipeline funcional e dockerizado.",
        "Arquitetura modular permite troca fácil de backbones.",
        "Próximos passos: Fine-tuning em dataset real (ex: Kaggle Fake News).",
        "Melhoria do OCR para fontes manuscritas ou distorcidas."
    ])

    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    try:
        import pptx
        generate_presentation("pbl_presentation.pptx")
    except ImportError:
        print("Erro: python-pptx não instalado. Execute: pip install python-pptx")

